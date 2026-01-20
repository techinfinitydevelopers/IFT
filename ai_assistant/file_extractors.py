"""
File Text Extraction Module for IFT Platform

Extracts text from uploaded documents (PDF, DOCX, PPT) and images (OCR)
for AI analysis. AI never receives raw files - only extracted text.
"""

import os
import logging
from io import BytesIO

logger = logging.getLogger(__name__)

# Maximum text length to prevent token overflow (approx 3000 tokens = 12000 chars)
MAX_EXTRACTED_TEXT_LENGTH = 12000
MAX_TEXT_PER_FILE = 4000


def extract_text_from_file(file_path, file_type, original_filename):
    """
    Extract text from uploaded file based on file type.
    
    Args:
        file_path: Path to the uploaded file
        file_type: Type of file ('document' or 'image')
        original_filename: Original filename for extension detection
        
    Returns:
        str: Extracted text (empty string if extraction fails)
    """
    try:
        extension = original_filename.lower().split('.')[-1] if '.' in original_filename else ''
        
        if file_type == 'document':
            if extension == 'pdf':
                return extract_from_pdf(file_path)
            elif extension in ('docx', 'doc'):
                return extract_from_docx(file_path)
            elif extension in ('pptx', 'ppt'):
                return extract_from_pptx(file_path)
            elif extension == 'txt':
                return extract_from_txt(file_path)
            else:
                logger.warning(f"Unsupported document type: {extension}")
                return ""
                
        elif file_type == 'image':
            return extract_from_image(file_path)
            
        else:
            # Video files - no text extraction
            return ""
            
    except Exception as e:
        logger.error(f"Error extracting text from {original_filename}: {str(e)}")
        return ""


def extract_from_pdf(file_path):
    """Extract text from PDF file using PyPDF2"""
    try:
        import PyPDF2
        
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                if len(''.join(text_parts)) > MAX_TEXT_PER_FILE:
                    text_parts.append(f"\n[... truncated, {len(reader.pages) - page_num} more pages ...]")
                    break
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        
        return clean_and_truncate_text('\n'.join(text_parts))
        
    except ImportError:
        logger.warning("PyPDF2 not installed. PDF extraction disabled.")
        return "[PDF text extraction unavailable]"
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        return ""


def extract_from_docx(file_path):
    """Extract text from DOCX file using python-docx"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if len('\n'.join(text_parts)) > MAX_TEXT_PER_FILE:
                text_parts.append("\n[... truncated ...]")
                break
            if para.text.strip():
                text_parts.append(para.text)
        
        return clean_and_truncate_text('\n'.join(text_parts))
        
    except ImportError:
        logger.warning("python-docx not installed. DOCX extraction disabled.")
        return "[DOCX text extraction unavailable]"
    except Exception as e:
        logger.error(f"DOCX extraction error: {str(e)}")
        return ""


def extract_from_pptx(file_path):
    """Extract text from PPTX file using python-pptx"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            if len('\n'.join(text_parts)) > MAX_TEXT_PER_FILE:
                text_parts.append(f"\n[... truncated, {len(prs.slides) - slide_num + 1} more slides ...]")
                break
                
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"Slide {slide_num}: {' '.join(slide_text)}")
        
        return clean_and_truncate_text('\n'.join(text_parts))
        
    except ImportError:
        logger.warning("python-pptx not installed. PPTX extraction disabled.")
        return "[PPTX text extraction unavailable]"
    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}")
        return ""


def extract_from_txt(file_path):
    """Extract text from plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return clean_and_truncate_text(text)
    except Exception as e:
        logger.error(f"TXT extraction error: {str(e)}")
        return ""


def extract_from_image(file_path):
    """Extract text from image using OCR (pytesseract)"""
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(file_path)
        
        # Convert to RGB if necessary
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Run OCR
        text = pytesseract.image_to_string(img, lang='eng')
        
        return clean_and_truncate_text(text)
        
    except ImportError:
        logger.warning("pytesseract not installed. Image OCR disabled.")
        return ""
    except Exception as e:
        logger.error(f"Image OCR error: {str(e)}")
        return ""


def clean_and_truncate_text(text):
    """Clean and truncate extracted text"""
    if not text:
        return ""
    
    # Split into lines
    lines = text.split('\n')
    
    # Join very short lines (likely word-by-word extraction)
    # If a line has fewer than 5 words and doesn't end with punctuation,
    # it's probably part of a sentence that was split
    cleaned_lines = []
    current_sentence = []
    
    for line in lines:
        line = line.strip()
        if not line:
            # Empty line - flush current sentence and add paragraph break
            if current_sentence:
                cleaned_lines.append(' '.join(current_sentence))
                current_sentence = []
            continue
        
        words = line.split()
        
        # If line is short and doesn't end with sentence-ending punctuation
        if len(words) <= 3 and not line[-1] in '.!?:':
            current_sentence.extend(words)
        else:
            # This looks like a complete line or end of sentence
            if current_sentence:
                current_sentence.extend(words)
                cleaned_lines.append(' '.join(current_sentence))
                current_sentence = []
            else:
                cleaned_lines.append(line)
    
    # Flush remaining
    if current_sentence:
        cleaned_lines.append(' '.join(current_sentence))
    
    # Join with proper paragraph breaks
    cleaned = '\n\n'.join(cleaned_lines)
    
    # Remove non-printable characters
    cleaned = ''.join(char for char in cleaned if char.isprintable() or char in '\n\t')
    
    # Truncate to max length
    if len(cleaned) > MAX_TEXT_PER_FILE:
        cleaned = cleaned[:MAX_TEXT_PER_FILE] + "\n[... text truncated for processing ...]"
    
    return cleaned


def combine_all_text(idea_text, extracted_texts):
    """
    Combine idea text with all extracted file texts.
    
    Args:
        idea_text: Main idea text from form fields
        extracted_texts: List of (filename, text) tuples
        
    Returns:
        str: Combined text ready for AI processing
    """
    combined = idea_text
    
    if extracted_texts:
        file_contents = []
        for filename, text in extracted_texts:
            if text and text.strip():
                file_contents.append(f"[From {filename}]:\n{text}")
        
        if file_contents:
            combined += "\n\n--- CONTENT FROM UPLOADED FILES ---\n" + "\n\n".join(file_contents)
    
    # Final truncation to prevent token overflow
    if len(combined) > MAX_EXTRACTED_TEXT_LENGTH:
        combined = combined[:MAX_EXTRACTED_TEXT_LENGTH] + "\n[... content truncated ...]"
    
    return combined


def process_uploaded_files(submission):
    """
    Process all uploaded files for a submission and extract text.
    Updates the extracted_text field on each UploadedFile.
    
    Args:
        submission: IdeaSubmission instance
        
    Returns:
        str: Combined extracted text from all files
    """
    from django.conf import settings
    
    extracted_texts = []
    
    for uploaded_file in submission.uploaded_files.all():
        try:
            # Get full file path
            file_path = os.path.join(settings.MEDIA_ROOT, str(uploaded_file.file))
            
            if os.path.exists(file_path):
                # Extract text
                text = extract_text_from_file(
                    file_path,
                    uploaded_file.file_type,
                    uploaded_file.original_filename
                )
                
                # Save extracted text to model
                if text:
                    uploaded_file.extracted_text = text
                    uploaded_file.save(update_fields=['extracted_text'])
                    extracted_texts.append((uploaded_file.original_filename, text))
                    
        except Exception as e:
            logger.error(f"Error processing file {uploaded_file.original_filename}: {str(e)}")
    
    return extracted_texts
